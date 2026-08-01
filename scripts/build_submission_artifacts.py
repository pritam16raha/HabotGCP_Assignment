#!/usr/bin/env python3
# Author: Pritam Raha <rahapritam32@gmail.com>
import hashlib
import json
import zipfile
from pathlib import Path

from openpyxl import Workbook, load_workbook
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUTPUT_DIRECTORY = ROOT / "submission"
PRESENTATION_PATH = OUTPUT_DIRECTORY / "Pritam_Raha_Engineering_Blueprint.pptx"
WORKBOOK_PATH = OUTPUT_DIRECTORY / "Pritam_Raha_Schema_Mapping.xlsx"
ARCHIVE_PATH = OUTPUT_DIRECTORY / "Pritam_Raha_Junior_Cloud_DevOps_Assignment.zip"
CHECKSUM_PATH = OUTPUT_DIRECTORY / "SHA256SUMS.txt"

NAVY = RGBColor(11, 24, 40)
NAVY_LIGHT = RGBColor(22, 42, 64)
TEAL = RGBColor(35, 201, 184)
BLUE = RGBColor(61, 136, 255)
AMBER = RGBColor(247, 183, 49)
RED = RGBColor(239, 93, 93)
WHITE = RGBColor(248, 250, 252)
MUTED = RGBColor(174, 189, 204)
INK = RGBColor(23, 37, 51)
PALE = RGBColor(232, 241, 246)


def _set_text(
    text_frame,
    text: str,
    *,
    size: int = 18,
    colour: RGBColor = INK,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
) -> None:
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = colour


def _text_box(slide, x, y, width, height, text, **text_options):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
    _set_text(box.text_frame, text, **text_options)
    return box


def _panel(slide, x, y, width, height, *, fill=WHITE, line=None, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    return shape


def _label(slide, x, y, width, text, *, fill=TEAL, colour=NAVY):
    shape = _panel(slide, x, y, width, 0.34, fill=fill)
    _set_text(
        shape.text_frame, text.upper(), size=10, colour=colour, bold=True, align=PP_ALIGN.CENTER
    )
    return shape


def _header(slide, number: int, title: str, kicker: str) -> None:
    _text_box(slide, 0.65, 0.35, 9.8, 0.3, kicker.upper(), size=10, colour=TEAL, bold=True)
    _text_box(slide, 0.65, 0.68, 11.7, 0.65, title, size=28, colour=WHITE, bold=True)
    _text_box(
        slide,
        11.9,
        0.36,
        0.75,
        0.3,
        f"{number:02d}",
        size=11,
        colour=MUTED,
        bold=True,
        align=PP_ALIGN.RIGHT,
    )


def _footer(slide) -> None:
    _text_box(
        slide,
        0.65,
        7.12,
        12.0,
        0.18,
        "PRITAM RAHA  •  rahapritam32@gmail.com  •  JUNIOR CLOUD AND DEVOPS ENGINEER",
        size=8,
        colour=MUTED,
    )


def _base_slide(presentation: Presentation, number: int, title: str, kicker: str):
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = NAVY
    _header(slide, number, title, kicker)
    _footer(slide)
    return slide


def _bullet_panel(slide, x, y, width, height, title, bullets, *, accent=TEAL):
    panel = _panel(slide, x, y, width, height, fill=NAVY_LIGHT, line=RGBColor(42, 65, 88))
    _text_box(
        slide, x + 0.28, y + 0.18, width - 0.56, 0.38, title, size=16, colour=WHITE, bold=True
    )
    text_box = slide.shapes.add_textbox(
        Inches(x + 0.28), Inches(y + 0.68), Inches(width - 0.56), Inches(height - 0.86)
    )
    frame = text_box.text_frame
    frame.clear()
    frame.word_wrap = True
    for index, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.space_after = Pt(9)
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(13)
        paragraph.font.color.rgb = MUTED
        paragraph.text = f"•  {bullet}"
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.06), Inches(height)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    return panel


def build_presentation() -> None:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY
    _label(slide, 0.75, 0.65, 2.0, "Hiring project")
    _text_box(
        slide,
        0.75,
        1.35,
        11.4,
        1.35,
        "Secure Student Onboarding\nDeployment & Automation Blueprint",
        size=31,
        colour=WHITE,
        bold=True,
    )
    _text_box(
        slide,
        0.78,
        3.1,
        10.5,
        0.72,
        "Google Cloud infrastructure  •  Fail-closed build gate  •  Deterministic data validation",
        size=17,
        colour=MUTED,
    )
    for index, (number, label) in enumerate(
        [("01", "Secure by default"), ("02", "Schema aligned"), ("03", "Objectively testable")]
    ):
        x = 0.78 + index * 3.45
        _panel(slide, x, 4.3, 3.05, 1.15, fill=NAVY_LIGHT, line=RGBColor(42, 65, 88))
        _text_box(slide, x + 0.22, 4.47, 0.62, 0.35, number, size=12, colour=TEAL, bold=True)
        _text_box(slide, x + 0.22, 4.86, 2.55, 0.36, label, size=15, colour=WHITE, bold=True)
    _text_box(
        slide,
        0.78,
        6.35,
        7.0,
        0.35,
        "Pritam Raha  •  rahapritam32@gmail.com",
        size=13,
        colour=WHITE,
        bold=True,
    )
    _text_box(slide, 0.78, 6.75, 7.0, 0.3, "Submission: 3 August 2026", size=11, colour=MUTED)

    slide = _base_slide(presentation, 2, "Incident converted into controls", "Problem framing")
    _bullet_panel(
        slide,
        0.65,
        1.52,
        3.85,
        4.95,
        "Observed failure",
        [
            "A raw application credential reached source control.",
            "A transactional schema change broke downstream analytics.",
            "Human memory was the only effective release gate.",
        ],
        accent=RED,
    )
    _bullet_panel(
        slide,
        4.73,
        1.52,
        3.85,
        4.95,
        "Engineering response",
        [
            "Remove key files: use narrow workload identities.",
            "Close every schema and validate it at three boundaries.",
            "Make a failed gate structurally unable to authorise release.",
        ],
        accent=AMBER,
    )
    _bullet_panel(
        slide,
        8.81,
        1.52,
        3.85,
        4.95,
        "Measurable outcome",
        [
            "Zero hardcoded credential findings.",
            "Twenty-seven fields agree across transport and warehouse.",
            "Forty-one infrastructure policy checks pass with zero failures.",
        ],
        accent=TEAL,
    )

    slide = _base_slide(presentation, 3, "End-to-end governed data flow", "Architecture")
    stages = [
        ("1", "Closed request", "Django REST\nFramework"),
        ("2", "Binary rules", "Six Yes or No\nchecks"),
        ("3", "D0 landing", "Immutable and\nencrypted"),
        ("4", "Schema event", "Pub/Sub\nApache Avro"),
        ("5", "D1 analytics", "BigQuery and\nrow filtering"),
    ]
    for index, (number, name, detail) in enumerate(stages):
        x = 0.62 + index * 2.52
        _panel(slide, x, 2.18, 2.05, 2.55, fill=WHITE)
        _label(slide, x + 0.17, 2.36, 0.45, number, fill=TEAL)
        _text_box(slide, x + 0.17, 2.9, 1.72, 0.48, name, size=15, colour=INK, bold=True)
        _text_box(slide, x + 0.17, 3.53, 1.72, 0.72, detail, size=12, colour=RGBColor(85, 103, 120))
        if index < len(stages) - 1:
            connector = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Inches(x + 2.05),
                Inches(3.45),
                Inches(x + 2.49),
                Inches(3.45),
            )
            connector.line.color.rgb = TEAL
            connector.line.width = Pt(2)
    _panel(slide, 1.05, 5.35, 11.25, 0.7, fill=NAVY_LIGHT, line=RGBColor(42, 65, 88))
    _text_box(
        slide,
        1.25,
        5.49,
        10.85,
        0.38,
        "Cloud Key Management Service encrypts storage, messaging and analytics  •  Every output retains source lineage",
        size=13,
        colour=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    slide = _base_slide(
        presentation, 4, "D0 raw landing: create, never mutate", "Terraform security"
    )
    controls = [
        ("PUBLIC", "Blocked", "Public access prevention is enforced"),
        ("WRITE", "Prefix only", "Object creation is conditioned to incoming/"),
        ("ENCRYPT", "Customer key", "Dedicated key rotates every 90 days"),
        ("RECOVER", "Layered", "Versions + retention + soft deletion"),
        ("AUDIT", "Recorded", "Access logs use a separate protected bucket"),
        ("DELETE", "Denied", "force_destroy is false"),
    ]
    for index, (tag, value, detail) in enumerate(controls):
        row, column = divmod(index, 3)
        x, y = 0.68 + column * 4.18, 1.55 + row * 2.32
        _panel(slide, x, y, 3.82, 1.85, fill=NAVY_LIGHT, line=RGBColor(42, 65, 88))
        _label(slide, x + 0.22, y + 0.18, 1.0, tag, fill=BLUE, colour=WHITE)
        _text_box(slide, x + 0.22, y + 0.65, 3.35, 0.4, value, size=18, colour=WHITE, bold=True)
        _text_box(slide, x + 0.22, y + 1.1, 3.35, 0.45, detail, size=11, colour=MUTED)

    slide = _base_slide(presentation, 5, "Least privilege is explicit", "Identity and access")
    identities = [
        ("Raw ingestor", "Create incoming objects", "No read • no list • no delete"),
        ("Data pipeline", "Read raw + maintain D1", "No identity or key administration"),
        ("Messaging agent", "Deliver valid events", "No human analytics access"),
        ("Dubai analyst", "Read Dubai rows only", "No other emirates • no writes"),
    ]
    for index, (name, allow, deny) in enumerate(identities):
        y = 1.55 + index * 1.22
        _panel(slide, 0.7, y, 11.95, 0.92, fill=NAVY_LIGHT, line=RGBColor(42, 65, 88))
        _text_box(slide, 0.98, y + 0.16, 2.3, 0.48, name, size=15, colour=WHITE, bold=True)
        _text_box(slide, 3.38, y + 0.16, 3.7, 0.48, allow, size=13, colour=TEAL, bold=True)
        _text_box(slide, 7.22, y + 0.16, 4.95, 0.48, deny, size=12, colour=MUTED)
    _text_box(
        slide,
        0.78,
        6.48,
        11.8,
        0.3,
        "No downloaded service-account key exists anywhere in the design.",
        size=14,
        colour=AMBER,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    slide = _base_slide(presentation, 6, "D1 staged and enforced", "BigQuery governance")
    _bullet_panel(
        slide,
        0.68,
        1.55,
        3.8,
        4.95,
        "Integrity",
        [
            "Twenty-seven required or repeated fields",
            "Unknown fields are never dropped",
            "Table deletion protection",
            "Ninety-six-hour time travel",
        ],
        accent=BLUE,
    )
    _bullet_panel(
        slide,
        4.76,
        1.55,
        3.8,
        4.95,
        "Performance",
        [
            "Daily ingestion partition",
            "Organisation and student clustering",
            "Ninety-day staging retention",
            "Regional managed service",
        ],
        accent=TEAL,
    )
    _bullet_panel(
        slide,
        8.84,
        1.55,
        3.8,
        4.95,
        "Row access",
        [
            "Pipeline policy: TRUE",
            "Analytics policy: emirate = DUBAI",
            "Dataset role still required",
            "Filtered role is policy-managed",
        ],
        accent=AMBER,
    )

    slide = _base_slide(
        presentation, 7, "One mapping; three enforced contracts", "Schema integrity"
    )
    _text_box(
        slide,
        0.8,
        1.45,
        11.75,
        0.55,
        "23 source leaves",
        size=20,
        colour=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    chevrons = [
        ("Serializer", "Types • limits • choices"),
        ("Canonical map", "One unique destination"),
        ("Pub/Sub", "Apache Avro schema"),
        ("BigQuery", "Type • mode • order"),
    ]
    for index, (name, detail) in enumerate(chevrons):
        x = 0.72 + index * 3.03
        shape = slide.shapes.add_shape(
            MSO_SHAPE.CHEVRON, Inches(x), Inches(2.38), Inches(2.8), Inches(1.55)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = TEAL if index in (0, 3) else BLUE
        shape.line.fill.background()
        _set_text(
            shape.text_frame,
            f"{name}\n{detail}",
            size=13,
            colour=WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
    _panel(slide, 2.1, 4.68, 9.1, 1.1, fill=NAVY_LIGHT, line=RGBColor(42, 65, 88))
    _text_box(
        slide,
        2.3,
        4.84,
        8.7,
        0.32,
        "+ 4 explicit system fields = 27 canonical fields",
        size=18,
        colour=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _text_box(
        slide,
        2.3,
        5.23,
        8.7,
        0.28,
        "Mismatch → publish or build fails; no silent data loss",
        size=12,
        colour=AMBER,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    slide = _base_slide(
        presentation, 8, "Six decisions. Only Yes or No.", "Deterministic business logic"
    )
    rules = [
        ("001", "Explicit processing consent"),
        ("002", "Consent chronology ≤ 30 days"),
        ("003", "Age from 3 through 21"),
        ("004", "Support decision is consistent"),
        ("005", "Diagnosis decision is consistent"),
        ("006", "Controlled lists are unique"),
    ]
    for index, (number, rule) in enumerate(rules):
        row, column = divmod(index, 2)
        x, y = 0.72 + column * 6.18, 1.5 + row * 1.62
        _panel(slide, x, y, 5.8, 1.25, fill=NAVY_LIGHT, line=RGBColor(42, 65, 88))
        _label(slide, x + 0.2, y + 0.18, 0.72, f"Rule {number}", fill=TEAL)
        _text_box(slide, x + 1.1, y + 0.24, 4.35, 0.63, rule, size=14, colour=WHITE, bold=True)
    _text_box(
        slide,
        1.2,
        6.43,
        10.9,
        0.34,
        "ACCEPT = Rule 001 ∧ Rule 002 ∧ Rule 003 ∧ Rule 004 ∧ Rule 005 ∧ Rule 006",
        size=14,
        colour=AMBER,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    slide = _base_slide(
        presentation, 9, "Serializer removes human interpretation", "Django REST Framework"
    )
    _bullet_panel(
        slide,
        0.68,
        1.55,
        5.75,
        4.95,
        "Field controls",
        [
            "Closed top-level and nested objects",
            "Exact minimum and maximum lengths",
            "Strict JSON booleans",
            "Case-sensitive controlled choices",
            "International phone and identifier patterns",
        ],
        accent=TEAL,
    )
    _bullet_panel(
        slide,
        6.72,
        1.55,
        5.92,
        4.95,
        "Cross-field controls",
        [
            "Submission freshness: 24 hours / 5-minute skew",
            "Consent interval: zero through 30 days",
            "Age evaluated on submission date",
            "Support and diagnosis consistency",
            "Database constraints preserve accepted invariants",
        ],
        accent=BLUE,
    )

    slide = _base_slide(presentation, 10, "A failure cannot reach release", "Poka-Yoke build graph")
    gates = [
        ("1", "Secret scan"),
        ("2", "Dependencies"),
        ("3", "Format + lint"),
        ("4", "Tests + contracts"),
        ("5", "Terraform + policy"),
    ]
    for index, (number, name) in enumerate(gates):
        x = 0.65 + index * 2.5
        _panel(slide, x, 2.05, 2.08, 1.32, fill=NAVY_LIGHT, line=RGBColor(42, 65, 88))
        _label(slide, x + 0.16, 2.22, 0.45, number, fill=BLUE, colour=WHITE)
        _text_box(
            slide,
            x + 0.16,
            2.67,
            1.72,
            0.4,
            name,
            size=13,
            colour=WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
    _text_box(
        slide,
        0.78,
        3.82,
        5.4,
        0.4,
        "ANY NON-ZERO EXIT",
        size=16,
        colour=RED,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _text_box(
        slide,
        7.1,
        3.82,
        5.4,
        0.4,
        "ALL GATES PASS",
        size=16,
        colour=TEAL,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _panel(slide, 1.25, 4.55, 4.4, 1.15, fill=RGBColor(67, 35, 43), line=RED)
    _text_box(
        slide,
        1.45,
        4.73,
        4.0,
        0.65,
        "QUARANTINE REPORT\nRelease job is structurally skipped",
        size=14,
        colour=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _panel(slide, 7.65, 4.55, 4.4, 1.15, fill=RGBColor(19, 65, 62), line=TEAL)
    _text_box(
        slide,
        7.85,
        4.73,
        4.0,
        0.65,
        "RELEASE AUTHORISATION\nImmutable commit recorded",
        size=14,
        colour=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    slide = _base_slide(
        presentation, 11, "Negative path is executable evidence", "Fail-closed demonstration"
    )
    _panel(slide, 0.72, 1.52, 7.75, 4.95, fill=RGBColor(8, 17, 29), line=RGBColor(42, 65, 88))
    code_lines = [
        "$ make demo-fail-closed",
        "",
        "Temporary insecure file created outside repository",
        "Credential-shaped literal detected",
        "Finding count: 1",
        "Release authorisation: NOT EMITTED",
        "Temporary file removed",
        "",
        "RESULT: PASSED",
    ]
    for index, line in enumerate(code_lines):
        colour = TEAL if line == "RESULT: PASSED" else (AMBER if line.startswith("$") else MUTED)
        _text_box(
            slide,
            1.05,
            1.82 + index * 0.46,
            7.0,
            0.32,
            line,
            size=12,
            colour=colour,
            bold=line in {"RESULT: PASSED", "$ make demo-fail-closed"},
        )
    _bullet_panel(
        slide,
        8.78,
        1.52,
        3.86,
        4.95,
        "Evidence quality",
        [
            "No real credential used",
            "No insecure fixture committed",
            "Finding evidence is redacted",
            "Test fails if detector misses input",
            "Quarantine path is failure-only",
        ],
        accent=TEAL,
    )

    slide = _base_slide(
        presentation, 12, "Validation result: ready to review", "Objective evidence"
    )
    metrics = [
        ("15", "automated tests", TEAL),
        ("27", "aligned fields", BLUE),
        ("41", "policy checks passed", TEAL),
        ("0", "policy failures", AMBER),
    ]
    for index, (value, label, colour) in enumerate(metrics):
        x = 0.7 + index * 3.08
        _panel(slide, x, 1.58, 2.72, 1.65, fill=NAVY_LIGHT, line=RGBColor(42, 65, 88))
        _text_box(
            slide,
            x + 0.15,
            1.77,
            2.42,
            0.62,
            value,
            size=29,
            colour=colour,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        _text_box(
            slide,
            x + 0.15,
            2.45,
            2.42,
            0.36,
            label,
            size=11,
            colour=MUTED,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
    _bullet_panel(
        slide,
        0.7,
        3.65,
        5.8,
        2.45,
        "Validated locally",
        [
            "Django and security tests",
            "Schema comparison",
            "Terraform plan test",
            "Checkov policy scan",
        ],
        accent=BLUE,
    )
    _bullet_panel(
        slide,
        6.82,
        3.65,
        5.8,
        2.45,
        "Deliberately not executed",
        [
            "No external project mutation",
            "No unapproved infrastructure apply",
            "No service-account key creation",
            "No production data used",
        ],
        accent=AMBER,
    )

    slide = _base_slide(
        presentation, 13, "Operate with discipline, improve with evidence", "Handoff"
    )
    _bullet_panel(
        slide,
        0.7,
        1.55,
        3.75,
        4.7,
        "Before deploy",
        [
            "Create protected remote state",
            "Review the exact plan",
            "Confirm data residency",
            "Approve with two-person review",
        ],
        accent=BLUE,
    )
    _bullet_panel(
        slide,
        4.77,
        1.55,
        3.75,
        4.7,
        "After deploy",
        [
            "Test prefix denial",
            "Test row filtering",
            "Publish valid and invalid events",
            "Verify zero drift",
        ],
        accent=TEAL,
    )
    _bullet_panel(
        slide,
        8.84,
        1.55,
        3.78,
        4.7,
        "Next improvement",
        [
            "Federated deployment identity",
            "Organisation policy constraints",
            "Regional recovery rehearsal",
            "Schema compatibility registry",
        ],
        accent=AMBER,
    )
    _text_box(
        slide,
        0.75,
        6.48,
        11.8,
        0.3,
        "Code, runbook, tests and wrapped schema workbook are included with this presentation.",
        size=13,
        colour=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    presentation.core_properties.title = (
        "Secure Student Onboarding Deployment and Automation Blueprint"
    )
    presentation.core_properties.author = "Pritam Raha <rahapritam32@gmail.com>"
    presentation.core_properties.subject = (
        "HabotConnect Junior Cloud and DevOps Engineer hiring project"
    )
    presentation.core_properties.keywords = (
        "Google Cloud, Terraform, Django REST Framework, fail-closed, BigQuery"
    )
    presentation.save(PRESENTATION_PATH)


def _style_worksheet(worksheet, widths: list[int]) -> None:
    dark_fill = PatternFill("solid", fgColor="0B1828")
    teal_fill = PatternFill("solid", fgColor="23C9B8")
    light_fill = PatternFill("solid", fgColor="E8F1F6")
    thin_border = Border(
        left=Side(style="thin", color="CBD5E1"),
        right=Side(style="thin", color="CBD5E1"),
        top=Side(style="thin", color="CBD5E1"),
        bottom=Side(style="thin", color="CBD5E1"),
    )
    worksheet.sheet_view.showGridLines = False
    worksheet.freeze_panes = "A2"
    worksheet.auto_filter.ref = worksheet.dimensions
    worksheet.page_setup.orientation = "landscape"
    worksheet.page_setup.fitToWidth = 1
    worksheet.sheet_properties.pageSetUpPr.fitToPage = True
    worksheet.oddFooter.center.text = "Pritam Raha | rahapritam32@gmail.com"
    worksheet.oddFooter.center.size = 9
    for column_index, width in enumerate(widths, start=1):
        worksheet.column_dimensions[get_column_letter(column_index)].width = width
    for row_index, row in enumerate(worksheet.iter_rows(), start=1):
        worksheet.row_dimensions[row_index].height = 36 if row_index == 1 else 54
        for cell in row:
            cell.alignment = Alignment(wrap_text=True, vertical="top")
            cell.border = thin_border
            cell.font = Font(name="Aptos", size=10, color="172533")
            if row_index == 1:
                cell.fill = dark_fill
                cell.font = Font(name="Aptos Display", size=11, bold=True, color="FFFFFF")
                cell.alignment = Alignment(wrap_text=True, vertical="center")
            elif row_index % 2 == 0:
                cell.fill = light_fill
    worksheet.sheet_properties.tabColor = "23C9B8"
    worksheet["A1"].fill = teal_fill
    worksheet["A1"].font = Font(name="Aptos Display", size=11, bold=True, color="0B1828")


def build_workbook() -> None:
    bigquery_fields = json.loads(
        (ROOT / "contracts/bigquery/student_onboarding.schema.json").read_text(encoding="utf-8")
    )
    mapping_details = {
        "schema_version": (
            "schema_version",
            "String",
            "Exactly 1.0",
            "Contract version must equal 1.0.",
        ),
        "submission_id": (
            "submission_id",
            "Universally unique identifier",
            "Exactly one valid identifier",
            "Required and globally unique.",
        ),
        "submitted_at": (
            "submitted_at",
            "Offset-aware timestamp",
            "No more than 24 hours old; no more than 5 minutes future",
            "Normalises to Coordinated Universal Time.",
        ),
        "organisation_id": (
            "organisation_id",
            "Universally unique identifier",
            "Exactly one valid identifier",
            "Required tenant boundary.",
        ),
        "student_external_id": (
            "student.external_id",
            "String",
            "6 to 32 characters; uppercase letters, digits, hyphens",
            "Pattern must match the school identifier rule.",
        ),
        "student_first_name": (
            "student.first_name",
            "String",
            "1 to 50 characters",
            "Letters with controlled separators only.",
        ),
        "student_last_name": (
            "student.last_name",
            "String",
            "1 to 50 characters",
            "Letters with controlled separators only.",
        ),
        "student_date_of_birth": (
            "student.date_of_birth",
            "Calendar date",
            "Completed age from 3 through 21 on submission date",
            "Deconstruction of Compliance into Yes or No rule 003.",
        ),
        "emirate": (
            "student.emirate",
            "Controlled string",
            "Exactly one of the seven United Arab Emirates emirates",
            "Case-sensitive code.",
        ),
        "school_name": (
            "student.school_name",
            "String",
            "2 to 120 characters",
            "Control characters are rejected.",
        ),
        "learning_support_required": (
            "student.learning_support_required",
            "JavaScript Object Notation boolean",
            "true or false only",
            "Deconstruction of Compliance into Yes or No rule 004.",
        ),
        "has_formal_diagnosis": (
            "student.has_formal_diagnosis",
            "JavaScript Object Notation boolean",
            "true or false only",
            "Deconstruction of Compliance into Yes or No rule 005.",
        ),
        "diagnosis_codes": (
            "student.diagnosis_codes",
            "List of controlled strings",
            "0 to 10 unique values; each 2 to 12 characters",
            "Required when formal diagnosis is true; empty when false.",
        ),
        "guardian_full_name": (
            "guardian.full_name",
            "String",
            "3 to 101 characters",
            "Letters with controlled separators only.",
        ),
        "guardian_email": (
            "guardian.email",
            "Electronic mail address",
            "Maximum 254 characters",
            "Django REST Framework electronic mail validation.",
        ),
        "guardian_phone_e164": (
            "guardian.phone_e164",
            "International telephone number",
            "9 to 16 characters including plus sign",
            "International Telecommunication Union E.164 pattern.",
        ),
        "guardian_relationship": (
            "guardian.relationship",
            "Controlled string",
            "Parent, legal guardian, or foster guardian",
            "Case-sensitive code.",
        ),
        "consent_to_process": (
            "guardian.consent_to_process",
            "JavaScript Object Notation boolean",
            "Must be true",
            "Deconstruction of Compliance into Yes or No rule 001.",
        ),
        "consent_timestamp": (
            "guardian.consent_timestamp",
            "Offset-aware timestamp",
            "0 to 30 days before submission",
            "Deconstruction of Compliance into Yes or No rule 002.",
        ),
        "support_areas": (
            "support.areas",
            "List of controlled strings",
            "0 to 8 unique values",
            "Populated only when support is required.",
        ),
        "requested_hours_per_week": (
            "support.requested_hours_per_week",
            "Whole number",
            "0 to 40 inclusive",
            "1 to 40 when support is true; zero when false.",
        ),
        "preferred_language": (
            "support.preferred_language",
            "Controlled string",
            "One of 6 supported languages",
            "Case-sensitive code.",
        ),
        "wheelchair_access_required": (
            "support.wheelchair_access_required",
            "JavaScript Object Notation boolean",
            "true or false only",
            "Strict binary value.",
        ),
        "dcyn_all_rules_passed": (
            "System generated",
            "Boolean",
            "Must be true for published events",
            "Conjunction of all six binary rules.",
        ),
        "dcyn_failed_rule_ids": (
            "System generated",
            "List of strings",
            "Empty for published events",
            "Failed data is quarantined before publication.",
        ),
        "ingested_at": (
            "System generated",
            "Offset-aware timestamp",
            "Set once at ingestion",
            "Normalises to Coordinated Universal Time.",
        ),
        "source_object_uri": (
            "System generated",
            "Cloud Storage uniform resource identifier",
            "Must include incoming/ and immutable generation",
            "Provides replay-safe source lineage.",
        ),
    }
    avro_fields = json.loads(
        (ROOT / "contracts/pubsub/student_onboarding.avsc").read_text(encoding="utf-8")
    )["fields"]
    avro_by_name = {field["name"]: field["type"] for field in avro_fields}

    workbook = Workbook()
    worksheet = workbook.active
    worksheet.title = "Field Mapping"
    worksheet.append(
        [
            "Source JavaScript Object Notation path",
            "Canonical field",
            "Input type",
            "Exact limit",
            "Validation decision",
            "Apache Avro type",
            "BigQuery type",
            "BigQuery mode",
            "BigQuery description",
        ]
    )
    for field in bigquery_fields:
        source, input_type, limit, decision = mapping_details[field["name"]]
        avro_type = avro_by_name[field["name"]]
        avro_text = (
            json.dumps(avro_type, separators=(",", ":"))
            if isinstance(avro_type, dict)
            else avro_type
        )
        worksheet.append(
            [
                source,
                field["name"],
                input_type,
                limit,
                decision,
                avro_text,
                field["type"],
                field["mode"],
                field["description"],
            ]
        )
    _style_worksheet(worksheet, [31, 30, 25, 39, 42, 34, 19, 19, 50])

    worksheet = workbook.create_sheet("Binary Logic")
    worksheet.append(["Rule identifier", "Binary question", "Yes condition", "No system action"])
    binary_rules = [
        (
            "Rule 001",
            "Did the guardian explicitly consent to processing?",
            "The value is the JavaScript Object Notation boolean true.",
            "Reject and quarantine the request.",
        ),
        (
            "Rule 002",
            "Was consent recorded no more than 30 days before submission?",
            "Consent is on or before submission and no more than 30 days old.",
            "Reject and return Rule 002.",
        ),
        (
            "Rule 003",
            "Is the student from 3 through 21 years old on the submission date?",
            "Completed age is inside the inclusive range.",
            "Reject and return Rule 003.",
        ),
        (
            "Rule 004",
            "Do support areas and weekly hours exactly match the support decision?",
            "Yes has at least one area and 1 to 40 hours; No has no areas and zero hours.",
            "Reject and return Rule 004.",
        ),
        (
            "Rule 005",
            "Do diagnosis codes exactly match the diagnosis decision?",
            "Yes has at least one code; No has no codes.",
            "Reject and return Rule 005.",
        ),
        (
            "Rule 006",
            "Does each controlled list contain unique values only?",
            "Diagnosis codes and support areas contain no duplicate.",
            "Reject and return Rule 006.",
        ),
    ]
    for row in binary_rules:
        worksheet.append(row)
    _style_worksheet(worksheet, [20, 54, 65, 38])

    worksheet = workbook.create_sheet("Access Control")
    worksheet.append(
        [
            "Identity",
            "Resource scope",
            "Allowed actions",
            "Denied by omission or condition",
            "Reason",
        ]
    )
    access_rows = [
        (
            "Raw ingestor service account",
            "Raw landing bucket incoming object prefix",
            "Create new objects",
            "Read, list, overwrite, delete, and all non-incoming prefixes",
            "Immutable least-privilege capture.",
        ),
        (
            "Data pipeline service account",
            "Raw incoming objects and staged dataset",
            "Read incoming objects, edit staged data, create query jobs",
            "Identity administration, key administration, and bucket administration",
            "Separates data processing from platform control.",
        ),
        (
            "Pub/Sub service agent",
            "Staged dataset and validated-events encryption key",
            "Deliver schema-valid messages and encrypt message data",
            "Human analytics access and infrastructure administration",
            "Managed delivery without a downloaded key.",
        ),
        (
            "Dubai analytics reader service account",
            "Staged dataset",
            "Run queries and read rows where emirate equals Dubai",
            "Other emirates and all write actions",
            "Dataset role plus BigQuery row access policy.",
        ),
        (
            "Cloud Storage service agent",
            "Raw and logging encryption key",
            "Encrypt and decrypt managed objects",
            "Key-ring administration and unrelated keys",
            "Narrow customer-managed encryption support.",
        ),
        (
            "BigQuery service agent",
            "Staged-data encryption key",
            "Encrypt and decrypt managed tables",
            "Key-ring administration and unrelated keys",
            "Narrow customer-managed encryption support.",
        ),
    ]
    for row in access_rows:
        worksheet.append(row)
    _style_worksheet(worksheet, [31, 38, 46, 49, 42])

    worksheet = workbook.create_sheet("Validation Evidence")
    worksheet.append(["Validation control", "Result", "Measured evidence", "Reproduction command"])
    evidence_rows = [
        ("Python formatting and lint", "Passed", "Zero findings", "make lint"),
        ("Bandit secure-code scan", "Passed", "Zero findings", "make lint"),
        (
            "Python dependency vulnerability audit",
            "Passed",
            "No known vulnerabilities",
            "pip-audit --strict --requirement requirements.txt",
        ),
        ("Django and pipeline tests", "Passed", "Thirteen tests", "make test"),
        ("Independent credential-scanner tests", "Passed", "Two tests", "make test"),
        (
            "Data contract comparison",
            "Passed",
            "Twenty-three input leaves plus four system fields equal twenty-seven canonical fields",
            "make contracts",
        ),
        ("Repository credential scan", "Passed", "Zero findings", "make contracts"),
        (
            "Synthetic fail-closed test",
            "Passed",
            "One unsafe temporary input blocked",
            "make demo-fail-closed",
        ),
        ("Terraform provider validation", "Passed", "Configuration valid", "make terraform"),
        ("Terraform mocked plan test", "Passed", "One plan test", "make terraform"),
        (
            "Checkov infrastructure security scan",
            "Passed",
            "Forty-one passed, zero failed, one documented terminal logging exception",
            "checkov --config-file .checkov.yml",
        ),
    ]
    for row in evidence_rows:
        worksheet.append(row)
    _style_worksheet(worksheet, [39, 18, 73, 43])

    workbook.properties.creator = "Pritam Raha <rahapritam32@gmail.com>"
    workbook.properties.title = "Student Onboarding Schema Mapping and Validation Evidence"
    workbook.properties.subject = "HabotConnect Junior Cloud and DevOps Engineer hiring project"
    workbook.save(WORKBOOK_PATH)

    verification_workbook = load_workbook(WORKBOOK_PATH)
    for verification_sheet in verification_workbook.worksheets:
        for row in verification_sheet.iter_rows():
            for cell in row:
                if cell.value is not None and cell.alignment.wrap_text is not True:
                    raise RuntimeError(
                        f"Wrap Text verification failed at {verification_sheet.title}!{cell.coordinate}"
                    )


def build_submission_archive() -> None:
    included_roots = [
        ROOT / ".github",
        ROOT / "architecture",
        ROOT / "backend",
        ROOT / "contracts",
        ROOT / "docs",
        ROOT / "examples",
        ROOT / "infrastructure",
        ROOT / "scripts",
        ROOT / "tests",
    ]
    included_files = [
        ROOT / ".checkov.yml",
        ROOT / ".editorconfig",
        ROOT / ".gitignore",
        ROOT / ".python-version",
        ROOT / "Makefile",
        ROOT / "README.md",
        ROOT / "pyproject.toml",
        ROOT / "requirements-artifacts.txt",
        ROOT / "requirements-dev.txt",
        ROOT / "requirements.txt",
        PRESENTATION_PATH,
        WORKBOOK_PATH,
    ]
    for included_root in included_roots:
        for path in included_root.rglob("*"):
            relative_parts = path.relative_to(ROOT).parts
            if (
                path.is_file()
                and ".terraform" not in relative_parts
                and "__pycache__" not in relative_parts
                and path.suffix not in {".pyc", ".tfplan", ".tfstate"}
            ):
                included_files.append(path)

    unique_files = sorted(set(included_files), key=lambda path: path.relative_to(ROOT).as_posix())
    checksum_lines = []
    for path in unique_files:
        digest = hashlib.sha256(path.read_bytes()).hexdigest()
        checksum_lines.append(f"{digest}  {path.relative_to(ROOT).as_posix()}")
    CHECKSUM_PATH.write_text("\n".join(checksum_lines) + "\n", encoding="utf-8")

    with zipfile.ZipFile(
        ARCHIVE_PATH, "w", compression=zipfile.ZIP_DEFLATED, compresslevel=9
    ) as archive:
        for path in [*unique_files, CHECKSUM_PATH]:
            archive.write(path, path.relative_to(ROOT).as_posix())

    with zipfile.ZipFile(ARCHIVE_PATH) as archive:
        corrupt_member = archive.testzip()
        if corrupt_member is not None:
            raise RuntimeError(f"Archive verification failed at {corrupt_member}")


def main() -> None:
    OUTPUT_DIRECTORY.mkdir(parents=True, exist_ok=True)
    build_presentation()
    build_workbook()
    build_submission_archive()
    print(f"Created {PRESENTATION_PATH.relative_to(ROOT)} (13 slides).")
    print(f"Created {WORKBOOK_PATH.relative_to(ROOT)} (4 worksheets; Wrap Text verified).")
    print(f"Created {ARCHIVE_PATH.relative_to(ROOT)} (verified upload package).")


if __name__ == "__main__":
    main()
